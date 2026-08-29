#!/usr/bin/env python3
"""
Smart India Hackathon 2026 — DarkTrace Nexus Presentation Generator
Professional dark-theme SIH presentation with 6 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ─── Color Palette ─────────────────────────────────────────────
BG_DARK       = RGBColor(0x0A, 0x0E, 0x1A)  # Deep navy background
BG_CARD       = RGBColor(0x10, 0x15, 0x25)  # Card background
BG_CARD_ALT   = RGBColor(0x14, 0x1A, 0x2E)  # Alternate card bg
BORDER_COLOR  = RGBColor(0x1E, 0x3A, 0x5F)  # Subtle border
CYAN          = RGBColor(0x00, 0xD4, 0xFF)  # Primary accent
CYAN_DIM      = RGBColor(0x00, 0x99, 0xBB)  # Dim cyan
BLUE          = RGBColor(0x3B, 0x82, 0xF6)  # Secondary accent
BLUE_DIM      = RGBColor(0x1E, 0x40, 0x80)
GREEN         = RGBColor(0x00, 0xE6, 0x96)  # Success / positive
RED           = RGBColor(0xFF, 0x45, 0x45)  # Warning
ORANGE        = RGBColor(0xFF, 0x99, 0x33)  # Orange accent
PURPLE        = RGBColor(0xA7, 0x8B, 0xFA)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM     = RGBColor(0xCC, 0xCC, 0xCC)
WHITE_MUTE    = RGBColor(0x88, 0x99, 0xAA)
GOLD          = RGBColor(0xFF, 0xD7, 0x00)

# ─── Presentation Setup ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=BG_DARK):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.space_after = Pt(line_spacing)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", spacing=4):
    """Add text with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(spacing)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=WHITE_DIM,
                    bullet_color=CYAN, font_name="Calibri", spacing=3, bold_items=False):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing)
        p.font.bold = bold_items
    return txBox

def add_flow_box(slide, left, top, width, height, text, fill=BG_CARD, border=CYAN,
                 text_color=WHITE, font_size=11, bold=False):
    """Add a flow diagram box."""
    shape = add_shape(slide, left, top, width, height, fill_color=fill, border_color=border, border_width=Pt(1.5))
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return shape

def add_arrow_down(slide, cx, top, length=Inches(0.25), color=CYAN):
    """Add a downward arrow connector."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - Inches(0.08), top, Inches(0.16), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_divider_line(slide, left, top, width, color=CYAN_DIM, thickness=Pt(1)):
    """Add a horizontal divider line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_glow_circle(slide, left, top, size, color=CYAN, opacity=40):
    """Add a decorative glow circle (semi-transparent)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_icon_badge(slide, left, top, width, height, icon_text, bg_color=CYAN, text_color=BG_DARK):
    """Add a small icon badge (emoji/icon in circle)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    p = shape.text_frame.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(height / Inches(1) * 14))
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return shape

def add_team_member_card(slide, left, top, width, height, name, role, index):
    """Add a team member card."""
    # Card background
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER_COLOR)
    # Avatar circle
    circle_size = Inches(0.55)
    circle_left = left + (width - circle_size) // 2
    circle_top = top + Inches(0.2)
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_size, circle_size)
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = CYAN
    avatar.line.fill.background()
    avatar.shadow.inherit = False
    # Initials
    initials = "".join([w[0] for w in name.split() if w])
    p = avatar.text_frame.paragraphs[0]
    p.text = initials[:2]
    p.font.size = Pt(13)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    try:
        avatar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    # Name
    add_text(slide, left + Inches(0.1), circle_top + circle_size + Inches(0.12),
             width - Inches(0.2), Inches(0.3), name,
             font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Role
    add_text(slide, left + Inches(0.1), circle_top + circle_size + Inches(0.42),
             width - Inches(0.2), Inches(0.5), role,
             font_size=9, color=CYAN_DIM, bold=False, alignment=PP_ALIGN.CENTER)

def add_stat_card(slide, left, top, width, height, label, value, accent_color=CYAN):
    """Add a statistics card."""
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=accent_color, border_width=Pt(1.5))
    add_text(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.3),
             value, font_size=24, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), Inches(0.3),
             label, font_size=10, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

def add_feature_card(slide, left, top, width, height, number, title, desc, accent=CYAN):
    """Add a feature card with number badge."""
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER_COLOR)
    # Number badge
    badge_size = Inches(0.35)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12), top + Inches(0.12), badge_size, badge_size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    badge.shadow.inherit = False
    p = badge.text_frame.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(11)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    try:
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    # Title
    add_text(slide, left + Inches(0.55), top + Inches(0.08), width - Inches(0.7), Inches(0.45),
             title, font_size=12.5, color=WHITE, bold=True)
    # Description
    add_text(slide, left + Inches(0.12), top + Inches(0.56), width - Inches(0.24), height - Inches(0.64),
             desc, font_size=10, color=WHITE_DIM)

def add_phase_box(slide, left, top, width, height, text, fill=BG_CARD_ALT, border=CYAN_DIM,
                  font_size=10, bold=False, text_color=WHITE_DIM):
    """Add a phase/process box."""
    shape = add_shape(slide, left, top, width, height, fill_color=fill, border_color=border)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    return shape


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)

# Decorative elements
add_glow_circle(slide1, Inches(-1), Inches(-1), Inches(3), CYAN)
add_glow_circle(slide1, Inches(11.5), Inches(5.5), Inches(2.5), BLUE)

# Top bar
add_shape(slide1, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

# SIH badge
sih_badge = add_shape(slide1, Inches(5.2), Inches(0.35), Inches(2.9), Inches(0.45),
                       fill_color=BG_CARD, border_color=CYAN, border_width=Pt(1.5))
p = sih_badge.text_frame.paragraphs[0]
p.text = "SIH 2026  |  NTRO  |  Cybersecurity"
p.font.size = Pt(11)
p.font.color.rgb = CYAN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
try:
    sih_badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass

# Main title
add_text(slide1, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.2),
         "DARKTRACE NEXUS", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri")

# Accent line under title
add_divider_line(slide1, Inches(4.5), Inches(3.05), Inches(4.3), CYAN, Pt(3))

# Subtitle
add_text(slide1, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
         "AI-Powered Dark Web Threat Actor Correlation & Attribution Platform",
         font_size=20, color=CYAN_DIM, bold=False, alignment=PP_ALIGN.CENTER)

# Tagline
add_text(slide1, Inches(2), Inches(4.0), Inches(9.3), Inches(0.4),
         "Defensive Cyber Threat Intelligence  |  Synthetic Data  |  Explainable Attribution",
         font_size=13, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

# Team members section
add_text(slide1, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.35),
         "TEAM DARKTRACE NEXUS", font_size=12, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri")

# Team member boxes
team_data = [
    ("D. PARDHIV", "Dark Web / OSINT Intelligence"),
    ("B. SREE KRISHNA GOWTHAM", "Team Lead & System Architect"),
    ("P. YUGANDHAR", "AI/ML & Behavioral Analysis"),
    ("C. SEETHA RAMADEVI", "Frontend, Backend & Visualization"),
    ("KVS. ANMOL", "Graph Analytics & Data Engineering"),
    ("A. SIRICHANDANA", "Cyber Threat Intelligence & Attribution"),
]

card_w = Inches(1.97)
card_h = Inches(1.0)
start_x = Inches(0.5)
gap = Inches(0.1)

for i, (name, role) in enumerate(team_data):
    x = start_x + i * (card_w + gap)
    y = Inches(5.15)
    card = add_shape(slide1, x, y, card_w, card_h, fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text(slide1, x + Inches(0.06), y + Inches(0.06), card_w - Inches(0.12), Inches(0.32),
             name, font_size=9.5, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide1, x + Inches(0.04), y + Inches(0.42), card_w - Inches(0.08), Inches(0.52),
             role, font_size=7, color=CYAN_DIM, alignment=PP_ALIGN.CENTER)

# Footer
add_text(slide1, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.3),
         "Problem Statement: SIH26151  |  Theme: Blockchain & Cybersecurity  |  github.com/degalapardhiv/DarkTraceNexus",
         font_size=9, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape(slide1, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=CYAN)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT & PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

# Top accent
add_shape(slide2, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

# Section header
add_text(slide2, Inches(0.6), Inches(0.3), Inches(5), Inches(0.4),
         "02  PROBLEM & SOLUTION", font_size=12, color=CYAN, bold=True)
add_divider_line(slide2, Inches(0.6), Inches(0.72), Inches(2.5), CYAN, Pt(2))

# ─── LEFT SIDE: The Problem ───
add_text(slide2, Inches(0.6), Inches(1.0), Inches(5.8), Inches(0.4),
         "THE PROBLEM", font_size=20, color=WHITE, bold=True)

add_text(slide2, Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.7),
         "Threat actors deliberately fragment their digital identity across multiple signals, making traditional correlation extremely difficult.",
         font_size=12, color=WHITE_DIM)

# Problem cards (fragmented signals)
problem_items = [
    ("Multiple Aliases", "Same actor, different names across forums"),
    ("Platform Identities", "Separate accounts on dark web marketplaces"),
    ("Cryptocurrency Wallets", "Different BTC/ETH addresses per operation"),
    ("PGP Identities", "Multiple encryption keys for authentication"),
    ("Infrastructure Shifts", "Changing domains, IPs, hosting providers"),
    ("Writing Style Variation", "Different linguistic patterns per persona"),
]

for i, (title, desc) in enumerate(problem_items):
    row = i // 2
    col = i % 2
    x = Inches(0.6) + col * Inches(3.05)
    y = Inches(2.3) + row * Inches(1.0)
    card = add_shape(slide2, x, y, Inches(2.85), Inches(0.85), fill_color=BG_CARD, border_color=BORDER_COLOR)
    # Red accent dot
    dot = slide2.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), y + Inches(0.15), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RED
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_text(slide2, x + Inches(0.3), y + Inches(0.08), Inches(2.4), Inches(0.3),
             title, font_size=11, color=WHITE, bold=True)
    add_text(slide2, x + Inches(0.3), y + Inches(0.38), Inches(2.4), Inches(0.4),
             desc, font_size=9, color=WHITE_MUTE)

# Result text
result_shape = add_shape(slide2, Inches(0.6), Inches(5.35), Inches(5.8), Inches(0.5),
                          fill_color=RGBColor(0x30, 0x15, 0x15), border_color=RED)
p = result_shape.text_frame.paragraphs[0]
p.text = "Result: Fragmented intelligence, missed correlations, slow investigations"
p.font.size = Pt(11)
p.font.color.rgb = RED
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
try:
    result_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass

# ─── CENTER: Arrow ───
arrow = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.55), Inches(3.2), Inches(0.45), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = CYAN
arrow.line.fill.background()
arrow.shadow.inherit = False

# ─── RIGHT SIDE: The Solution ───
add_text(slide2, Inches(7.2), Inches(1.0), Inches(5.8), Inches(0.4),
         "DARKTRACE NEXUS SOLUTION", font_size=20, color=GREEN, bold=True)

add_text(slide2, Inches(7.2), Inches(1.45), Inches(5.8), Inches(0.7),
         "Multi-signal correlation engine that unifies fragmented intelligence into explainable attribution.",
         font_size=12, color=WHITE_DIM)

# Solution flow
solution_flow = [
    ("Fragmented Threat Intelligence", WHITE_MUTE),
    ("Multi-Signal Correlation Engine", CYAN),
    ("Graph-Based Relationship Analysis", BLUE),
    ("Explainable Attribution Scoring", GREEN),
    ("Analyst Investigation & Reporting", WHITE),
]

for i, (label, color) in enumerate(solution_flow):
    y = Inches(2.3) + i * Inches(0.72)
    box_w = Inches(5.3)
    box_h = Inches(0.52)
    x = Inches(7.2)

    if i == 0:
        fill = RGBColor(0x25, 0x15, 0x15)
        border = RED
    elif i == len(solution_flow) - 1:
        fill = RGBColor(0x0F, 0x25, 0x18)
        border = GREEN
    else:
        fill = BG_CARD
        border = color

    phase = add_phase_box(slide2, x, y, box_w, box_h, label, fill=fill, border=border,
                          font_size=12, bold=True, text_color=color)

    if i < len(solution_flow) - 1:
        add_arrow_down(slide2, x + box_w // 2, y + box_h, Inches(0.18), CYAN_DIM)

# Key capabilities highlight
add_text(slide2, Inches(7.2), Inches(6.0), Inches(5.8), Inches(0.3),
         "Core Capabilities:", font_size=10, color=CYAN, bold=True)

caps = ["OSINT Correlation", "Behavioral Analysis", "Stylometry", "Infrastructure Mapping", "Evidence Chain", "Explainable AI"]
for i, cap in enumerate(caps):
    col = i % 3
    row = i // 3
    x = Inches(7.2) + col * Inches(1.85)
    y = Inches(6.3) + row * Inches(0.35)
    add_text(slide2, x, y, Inches(1.8), Inches(0.3),
             f"  {cap}", font_size=9, color=WHITE_DIM)

# Speaking note placeholder (bottom)
add_text(slide2, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
         "Speaking Note: Threat actors fragment identities across aliases, wallets, PGP keys, and infrastructure. DarkTrace Nexus correlates these signals through AI analysis and graph intelligence to deliver explainable attribution for analysts.",
         font_size=8, color=RGBColor(0x55, 0x66, 0x77), alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

# Top accent
add_shape(slide3, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

add_text(slide3, Inches(0.6), Inches(0.3), Inches(5), Inches(0.4),
         "03  SYSTEM ARCHITECTURE", font_size=12, color=CYAN, bold=True)
add_divider_line(slide3, Inches(0.6), Inches(0.72), Inches(2.5), CYAN, Pt(2))

add_text(slide3, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "End-to-End Intelligence Pipeline", font_size=22, color=WHITE, bold=True)

# Subtle background panel unifying the architecture diagram
add_shape(slide3, Inches(0.3), Inches(1.3), Inches(12.73), Inches(5.05),
          fill_color=RGBColor(0x0D, 0x13, 0x24), border_color=RGBColor(0x16, 0x2A, 0x47), border_width=Pt(1))

# ─── LEFT COLUMN: Data Sources ───
col1_x = Inches(0.4)
col1_w = Inches(2.6)

add_text(slide3, col1_x, Inches(1.45), col1_w, Inches(0.3),
         "DATA / INTELLIGENCE SOURCES", font_size=10, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

sources = [
    "Dark Web Forums",
    "Paste Sites & Pastebin",
    "Marketplace Listings",
    "PGP Key Servers",
    "Blockchain Transactions",
    "WHOIS / DNS Records",
    "OSINT Feeds",
]

src_y = Inches(1.8)
for i, src in enumerate(sources):
    y = src_y + i * Inches(0.45)
    box = add_phase_box(slide3, col1_x, y, col1_w, Inches(0.38), src,
                        fill=BG_CARD, border=BORDER_COLOR, font_size=9, text_color=WHITE_DIM)

# Arrow from sources to collection
arrow1 = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.15), Inches(1.9), Inches(0.35), Inches(0.4))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = CYAN
arrow1.line.fill.background()
arrow1.shadow.inherit = False

# ─── MIDDLE COLUMN: Processing Pipeline ───
col2_x = Inches(3.6)
col2_w = Inches(3.0)

pipeline_steps = [
    ("Collection & Normalization", BG_CARD_ALT, CYAN_DIM),
    ("Entity Extraction", BG_CARD_ALT, CYAN_DIM),
    ("Correlation Engine", BG_CARD, CYAN),
]

for i, (label, fill, border) in enumerate(pipeline_steps):
    y = Inches(1.5) + i * Inches(0.72)
    box = add_phase_box(slide3, col2_x, y, col2_w, Inches(0.52), label,
                        fill=fill, border=border, font_size=11, bold=True, text_color=WHITE)
    if i < len(pipeline_steps) - 1:
        add_arrow_down(slide3, col2_x + col2_w // 2, y + Inches(0.52), Inches(0.18), CYAN_DIM)

# Connector: Correlation Engine -> AI/ML Analysis box
add_arrow_down(slide3, col2_x + col2_w // 2, Inches(3.46), Inches(0.39), CYAN)

# AI/ML Analysis Box
ai_y = Inches(3.85)
ai_box = add_shape(slide3, col2_x, ai_y, col2_w, Inches(2.0),
                    fill_color=RGBColor(0x0D, 0x1F, 0x35), border_color=CYAN, border_width=Pt(2))
add_text(slide3, col2_x + Inches(0.1), ai_y + Inches(0.05), col2_w - Inches(0.2), Inches(0.25),
         "AI/ML ANALYSIS", font_size=10, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

ai_items = [
    "Behavioral Analysis",
    "Stylometric Profiling",
    "Entity Similarity",
    "Temporal Correlation",
    "8-Factor Attribution",
]
for i, item in enumerate(ai_items):
    add_text(slide3, col2_x + Inches(0.25), ai_y + Inches(0.35) + i * Inches(0.3),
             col2_w - Inches(0.4), Inches(0.25),
             f"  {item}", font_size=9, color=WHITE_DIM)

# Arrow to graph
arrow2 = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.75), Inches(3.5), Inches(0.35), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = CYAN
arrow2.line.fill.background()
arrow2.shadow.inherit = False

# ─── RIGHT COLUMN: Graph & Output ───
col3_x = Inches(7.2)
col3_w = Inches(2.8)

# Graph Analytics box
graph_box = add_shape(slide3, col3_x, Inches(1.5), col3_w, Inches(2.2),
                       fill_color=RGBColor(0x0D, 0x1F, 0x35), border_color=BLUE, border_width=Pt(2))
add_text(slide3, col3_x + Inches(0.1), Inches(1.55), col3_w - Inches(0.2), Inches(0.25),
         "GRAPH ANALYTICS", font_size=10, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

graph_entities = ["Actor", "Alias", "PGP Key", "Wallet", "Infrastructure", "Evidence"]
for i, ent in enumerate(graph_entities):
    row = i // 3
    col = i % 3
    ex = col3_x + Inches(0.15) + col * Inches(0.85)
    ey = Inches(1.9) + row * Inches(0.55)
    entity_box = add_shape(slide3, ex, ey, Inches(0.78), Inches(0.4),
                           fill_color=BG_CARD, border_color=BLUE_DIM)
    p = entity_box.text_frame.paragraphs[0]
    p.text = ent
    p.font.size = Pt(8)
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    try:
        entity_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass

# Connection lines description
add_text(slide3, col3_x + Inches(0.1), Inches(2.95), col3_w - Inches(0.2), Inches(0.6),
         "Actor ↔ Alias ↔ PGP ↔ Wallet ↔ Infra ↔ Evidence",
         font_size=8, color=BLUE_DIM, alignment=PP_ALIGN.CENTER)

# Arrow to attribution
arrow3 = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8.6), Inches(3.8), Inches(0.2), Inches(0.3))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = CYAN
arrow3.line.fill.background()
arrow3.shadow.inherit = False

# Attribution Engine
attr_box = add_shape(slide3, col3_x, Inches(4.2), col3_w, Inches(0.5),
                      fill_color=RGBColor(0x0F, 0x25, 0x18), border_color=GREEN, border_width=Pt(2))
p = attr_box.text_frame.paragraphs[0]
p.text = "Explainable Attribution Engine"
p.font.size = Pt(10)
p.font.color.rgb = GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
try:
    attr_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass

# Arrow: Attribution Engine -> Tech Stack (implementation layer)
arrow4 = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.0), Inches(4.4), Inches(0.3), Inches(0.2))
arrow4.fill.solid()
arrow4.fill.fore_color.rgb = CYAN
arrow4.line.fill.background()
arrow4.shadow.inherit = False

# ─── FAR RIGHT: Tech Stack ───
far_x = Inches(10.3)
far_w = Inches(2.8)

add_text(slide3, far_x, Inches(1.5), far_w, Inches(0.3),
         "TECHNOLOGY STACK", font_size=10, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

tech_layers = [
    ("Next.js 14", "Frontend Dashboard", CYAN),
    ("FastAPI", "REST API Backend", BLUE),
    ("PostgreSQL 16", "Primary Database", GREEN),
    ("Neo4j 5", "Graph Analytics", PURPLE),
    ("Python", "AI/ML Analysis", ORANGE),
    ("SSE / Realtime", "Live Dashboard", CYAN_DIM),
]

for i, (tech, desc, color) in enumerate(tech_layers):
    y = Inches(1.85) + i * Inches(0.62)
    card = add_shape(slide3, far_x, y, far_w, Inches(0.52), fill_color=BG_CARD, border_color=color, border_width=Pt(1.5))
    add_text(slide3, far_x + Inches(0.1), y + Inches(0.02), far_w - Inches(0.2), Inches(0.25),
             tech, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide3, far_x + Inches(0.1), y + Inches(0.27), far_w - Inches(0.2), Inches(0.2),
             desc, font_size=8, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

# ─── Bottom: Data flow summary ───
flow_labels = [
    ("Sources", Inches(0.5)),
    ("Collection", Inches(3.2)),
    ("Analysis", Inches(5.0)),
    ("Graph", Inches(7.2)),
    ("Attribution", Inches(9.0)),
    ("Dashboard", Inches(11.0)),
]

add_divider_line(slide3, Inches(0.4), Inches(6.45), Inches(12.5), BORDER_COLOR, Pt(1))

for label, x in flow_labels:
    add_text(slide3, x, Inches(6.55), Inches(1.5), Inches(0.25),
             label, font_size=8, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

# Speaking note
add_text(slide3, Inches(0.6), Inches(7.0), Inches(12), Inches(0.3),
         "Speaking Note: Our architecture follows a pipeline from data collection through entity extraction, AI-powered correlation, graph analytics, and explainable attribution — all surfaced through a real-time Next.js dashboard backed by FastAPI and PostgreSQL.",
         font_size=8, color=RGBColor(0x55, 0x66, 0x77))


# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — INNOVATION & KEY FEATURES
# ═══════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_shape(slide4, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

add_text(slide4, Inches(0.6), Inches(0.3), Inches(5), Inches(0.4),
         "04  INNOVATION & KEY FEATURES", font_size=12, color=CYAN, bold=True)
add_divider_line(slide4, Inches(0.6), Inches(0.72), Inches(2.5), CYAN, Pt(2))

add_text(slide4, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "Six Core Differentiators", font_size=22, color=WHITE, bold=True)

# Feature cards (2 rows x 3 columns)
features = [
    (1, "Multi-Signal Actor Correlation",
     "Combines aliases, PGP keys, cryptocurrency wallets, infrastructure, behavioral patterns, and stylometric analysis into unified actor profiles.",
     CYAN),
    (2, "Explainable Attribution",
     "8-factor weighted scoring shows WHY actors are correlated — not just an AI score. Analysts see the reasoning behind every attribution.",
     GREEN),
    (3, "Interactive Relationship Graph",
     "ReactFlow-powered graph visualizes Actor → Alias → PGP → Wallet → Infrastructure → Evidence relationships with confidence thresholds.",
     BLUE),
    (4, "Behavioral & Stylometric Intelligence",
     "Identifies similarities across posting patterns, activity timing, vocabulary richness, sentence structure, and linguistic fingerprints.",
     ORANGE),
    (5, "Evidence Chain Integrity",
     "SHA-256 evidence hashing, timestamped records, and linked intelligence entities create verifiable, tamper-evident investigation trails.",
     PURPLE),
    (6, "Investigation Workspace",
     "End-to-end workflow: Profile → Graph → Attribution → Evidence → Timeline → Report generation with JSON/CSV export.",
     CYAN),
]

card_w = Inches(3.9)
card_h = Inches(1.6)
gap_x = Inches(0.2)
gap_y = Inches(0.15)
start_x = Inches(0.5)
start_y = Inches(1.5)

for i, (num, title, desc, accent) in enumerate(features):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_feature_card(slide4, x, y, card_w, card_h, num, title, desc, accent)

# ─── Attribution Concept Diagram (bottom right) ───
add_text(slide4, Inches(0.6), Inches(4.95), Inches(5), Inches(0.3),
         "EXPLAINABLE ATTRIBUTION CONCEPT", font_size=10, color=CYAN, bold=True)

# Attribution factors
attr_factors = [
    ("Identity", CYAN),
    ("Behavior", BLUE),
    ("Stylometry", GREEN),
    ("Infrastructure", ORANGE),
    ("PGP", PURPLE),
    ("Wallet", CYAN_DIM),
    ("Temporal", WHITE_MUTE),
    ("Source", WHITE_DIM),
]

factor_x = Inches(0.5)
factor_y = Inches(5.35)
factor_w = Inches(1.1)
factor_h = Inches(0.4)
factor_gap = Inches(0.15)

for i, (label, color) in enumerate(attr_factors):
    x = factor_x + (i % 4) * (factor_w + factor_gap)
    y = factor_y + (i // 4) * (factor_h + Inches(0.1))
    box = add_shape(slide4, x, y, factor_w, factor_h, fill_color=BG_CARD, border_color=color)
    p = box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    try:
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass

# Arrow to weighted
arrow_attr = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.5), Inches(6.4), Inches(0.2), Inches(0.25))
arrow_attr.fill.solid()
arrow_attr.fill.fore_color.rgb = CYAN
arrow_attr.line.fill.background()
arrow_attr.shadow.inherit = False

# Weighted result
weight_box = add_shape(slide4, Inches(1.0), Inches(6.7), Inches(3.0), Inches(0.4),
                        fill_color=RGBColor(0x0F, 0x25, 0x18), border_color=GREEN, border_width=Pt(2))
p = weight_box.text_frame.paragraphs[0]
p.text = "Weighted Correlation → Confidence Assessment"
p.font.size = Pt(10)
p.font.color.rgb = GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
try:
    weight_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass

# ─── Key metrics highlight (right side) ───
add_text(slide4, Inches(5.5), Inches(4.95), Inches(7), Inches(0.3),
         "PLATFORM CAPABILITIES", font_size=10, color=CYAN, bold=True)

capabilities = [
    "19 API endpoints with full CRUD operations",
    "18 database tables with async SQLAlchemy ORM",
    "PostgreSQL 16 + Neo4j 5 dual persistence",
    "8-factor weighted attribution scoring engine",
    "Real-time SSE dashboard with heartbeat",
    "JWT authentication with Argon2 + RBAC",
    "Interactive ReactFlow relationship graph",
    "SHA-256 evidence integrity verification",
    "JSON/CSV export for all intelligence data",
    "Docker multi-stage builds + cloud deployment",
]

for i, cap in enumerate(capabilities):
    row = i // 2
    col = i % 2
    x = Inches(5.5) + col * Inches(3.7)
    y = Inches(5.35) + row * Inches(0.35)
    add_text(slide4, x, y, Inches(3.6), Inches(0.3),
             f"  {cap}", font_size=9, color=WHITE_DIM)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — FEASIBILITY, IMPACT & DEMONSTRATION
# ═══════════════════════════════════════════════════════════════

slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_shape(slide5, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

add_text(slide5, Inches(0.6), Inches(0.3), Inches(8), Inches(0.4),
         "05  FEASIBILITY, IMPACT & DEMONSTRATION", font_size=12, color=CYAN, bold=True)
add_divider_line(slide5, Inches(0.6), Inches(0.72), Inches(2.5), CYAN, Pt(2))

add_text(slide5, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "Working Prototype — Not Just a Concept", font_size=22, color=WHITE, bold=True)

# Subtle column panels to group the three sections
add_shape(slide5, Inches(0.45), Inches(1.3), Inches(5.75), Inches(4.95),
          fill_color=RGBColor(0x0D, 0x13, 0x24), border_color=RGBColor(0x16, 0x2A, 0x47), border_width=Pt(1))
add_shape(slide5, Inches(9.85), Inches(1.3), Inches(3.1), Inches(5.8),
          fill_color=RGBColor(0x0D, 0x13, 0x24), border_color=RGBColor(0x16, 0x2A, 0x47), border_width=Pt(1))

# ─── LEFT: Working Prototype Stats ───
add_text(slide5, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.3),
         "WORKING PROTOTYPE", font_size=12, color=GREEN, bold=True)

stats = [
    ("50", "Synthetic Threat Actors"),
    ("10", "Cross-Actor Correlations"),
    ("521", "Evidence Records"),
    ("130", "Timeline Events"),
    ("1,282", "Intelligence Posts"),
    ("19", "API Endpoints"),
]

stat_w = Inches(1.7)
stat_h = Inches(1.1)
stat_gap = Inches(0.12)

for i, (value, label) in enumerate(stats):
    row = i // 3
    col = i % 3
    x = Inches(0.6) + col * (stat_w + stat_gap)
    y = Inches(1.85) + row * (stat_h + Inches(0.1))
    add_stat_card(slide5, x, y, stat_w, stat_h, label, value, CYAN)

# ─── Working Features ───
add_text(slide5, Inches(0.6), Inches(4.3), Inches(5.5), Inches(0.3),
         "VERIFIED WORKING FEATURES", font_size=12, color=CYAN, bold=True)

working_features = [
    "Interactive relationship graph (ReactFlow)",
    "Explainable attribution with 8-factor scoring",
    "Investigation workspace with step-by-step workflow",
    "Report generation with JSON/CSV export",
    "Real-time SSE dashboard updates",
    "Global cross-entity search",
    "Timeline analysis with event filtering",
    "Evidence chain with SHA-256 verification",
]

for i, feat in enumerate(working_features):
    row = i // 2
    col = i % 2
    x = Inches(0.6) + col * Inches(3.0)
    y = Inches(4.65) + row * Inches(0.35)
    add_text(slide5, x, y, Inches(2.9), Inches(0.3),
             f"  {feat}", font_size=9, color=WHITE_DIM)

# ─── CENTER: Demo Workflow ───
# Container panel behind demo workflow (created first so header/steps render on top)
demo_panel = add_shape(slide5, Inches(6.6), Inches(1.4), Inches(3.25), Inches(5.2),
                       fill_color=BG_CARD, border_color=BORDER_COLOR, border_width=Pt(1))

add_text(slide5, Inches(6.6), Inches(1.5), Inches(3), Inches(0.3),
         "DEMO WORKFLOW", font_size=12, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

demo_steps = [
    "Dashboard Overview",
    "Threat Actor List",
    "Actor Profile Detail",
    "Relationship Graph",
    "Attribution Analysis",
    "Evidence Chain",
    "Timeline View",
    "Investigation Report",
]

demo_x = Inches(6.7)
demo_w = Inches(2.8)
for i, step in enumerate(demo_steps):
    y = Inches(1.9) + i * Inches(0.58)
    box = add_phase_box(slide5, demo_x, y, demo_w, Inches(0.42), step,
                        fill=BG_CARD, border=CYAN_DIM if i < len(demo_steps)-1 else GREEN,
                        font_size=10, bold=(i==0 or i==len(demo_steps)-1),
                        text_color=WHITE if i==0 or i==len(demo_steps)-1 else WHITE_DIM)
    if i < len(demo_steps) - 1:
        add_arrow_down(slide5, demo_x + demo_w // 2, y + Inches(0.42), Inches(0.14), CYAN_DIM)

# ─── RIGHT: Impact ───
add_text(slide5, Inches(10.0), Inches(1.5), Inches(3), Inches(0.3),
         "IMPACT", font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

impacts = [
    ("Reduces fragmented CTI analysis", GREEN),
    ("Discovers hidden actor relationships", CYAN),
    ("Provides evidence-backed correlation", BLUE),
    ("Improves investigation efficiency", ORANGE),
    ("Supports transparent attribution", PURPLE),
    ("Scalable CTI/SOC foundation", WHITE),
]

for i, (impact, color) in enumerate(impacts):
    y = Inches(1.9) + i * Inches(0.5)
    dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.1), y + Inches(0.08), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_text(slide5, Inches(10.3), y, Inches(2.8), Inches(0.4),
             impact, font_size=9, color=WHITE_DIM)

# ─── Bottom: Feasibility/Scalability ───
add_text(slide5, Inches(10.0), Inches(4.9), Inches(3), Inches(0.3),
         "SCALABILITY", font_size=12, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

scalability = [
    "Modular microservice architecture",
    "API-driven design (REST + SSE)",
    "Graph-based relationship model",
    "Managed database compatibility",
    "Cloud-deployment ready",
    "Docker containerization",
]

for i, item in enumerate(scalability):
    y = Inches(5.25) + i * Inches(0.3)
    add_text(slide5, Inches(10.0), y, Inches(3), Inches(0.28),
             f"  {item}", font_size=9, color=WHITE_DIM)


# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — TEAM, FUTURE SCOPE & CLOSING
# ═══════════════════════════════════════════════════════════════

slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)

add_shape(slide6, Inches(0), Inches(0), W, Inches(0.06), fill_color=CYAN)

add_text(slide6, Inches(0.6), Inches(0.3), Inches(5), Inches(0.4),
         "06  TEAM & FUTURE SCOPE", font_size=12, color=CYAN, bold=True)
add_divider_line(slide6, Inches(0.6), Inches(0.72), Inches(2.5), CYAN, Pt(2))

add_text(slide6, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "Meet the Team Behind DarkTrace Nexus", font_size=22, color=WHITE, bold=True)

# ─── Team Members Grid (2 rows x 3 cols) ───
team_full = [
    ("D. PARDHIV", "Dark Web / OSINT Intelligence"),
    ("B. SREE KRISHNA GOWTHAM", "Team Lead & System Architect"),
    ("P. YUGANDHAR", "AI/ML & Behavioral Analysis"),
    ("C. SEETHA RAMADEVI", "Frontend, Backend & Visualization"),
    ("KVS. ANMOL", "Graph Analytics & Data Engineering"),
    ("A. SIRICHANDANA", "Cyber Threat Intelligence & Attribution"),
]

card_w = Inches(3.9)
card_h = Inches(1.45)
gap_x = Inches(0.2)
gap_y = Inches(0.15)
start_x = Inches(0.5)
start_y = Inches(1.4)

for i, (name, role) in enumerate(team_full):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_shape(slide6, x, y, card_w, card_h, fill_color=BG_CARD, border_color=BORDER_COLOR)

    # Avatar circle
    circle_size = Inches(0.5)
    avatar = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), circle_size, circle_size)
    avatar.fill.solid()
    colors = [CYAN, BLUE, GREEN, ORANGE, PURPLE, CYAN_DIM]
    avatar.fill.fore_color.rgb = colors[i]
    avatar.line.fill.background()
    avatar.shadow.inherit = False
    initials = "".join([w[0] for w in name.split() if w])
    p = avatar.text_frame.paragraphs[0]
    p.text = initials[:2]
    p.font.size = Pt(12)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    try:
        avatar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass

    # Name
    add_text(slide6, x + Inches(0.75), y + Inches(0.1), card_w - Inches(0.9), Inches(0.3),
             name, font_size=12, color=WHITE, bold=True)
    # Role
    add_text(slide6, x + Inches(0.75), y + Inches(0.4), card_w - Inches(0.9), Inches(0.45),
             role, font_size=9, color=colors[i], bold=True)

# ─── Future Scope ───
add_text(slide6, Inches(0.6), Inches(4.55), Inches(6), Inches(0.3),
         "FUTURE SCOPE", font_size=12, color=CYAN, bold=True)
add_divider_line(slide6, Inches(0.6), Inches(4.9), Inches(1.5), CYAN, Pt(1))

future_items_left = [
    "Larger intelligence datasets & real-time ingestion",
    "Advanced entity-resolution models",
    "Improved multilingual stylometry",
    "Automated CTI enrichment pipelines",
]

future_items_right = [
    "Advanced graph anomaly detection",
    "SIEM/SOC platform integration",
    "Threat intelligence platform connectors",
    "Scalable realtime intelligence pipelines",
]

for i, item in enumerate(future_items_left):
    y = Inches(5.05) + i * Inches(0.35)
    add_text(slide6, Inches(0.6), y, Inches(5.8), Inches(0.3),
             f"  {item}", font_size=10, color=WHITE_DIM)

for i, item in enumerate(future_items_right):
    y = Inches(5.05) + i * Inches(0.35)
    add_text(slide6, Inches(6.6), y, Inches(6), Inches(0.3),
             f"  {item}", font_size=10, color=WHITE_DIM)

# ─── GitHub & Closing ───
add_divider_line(slide6, Inches(0.6), Inches(6.5), Inches(12), BORDER_COLOR, Pt(1))

# GitHub
add_text(slide6, Inches(0.6), Inches(6.62), Inches(4.6), Inches(0.3),
         "github.com/degalapardhiv/DarkTraceNexus", font_size=11, color=CYAN, bold=True)

# Closing statement
closing_box = add_shape(slide6, Inches(5.5), Inches(6.55), Inches(7.3), Inches(0.5),
                         fill_color=RGBColor(0x0D, 0x1F, 0x35), border_color=CYAN, border_width=Pt(1.5))
p = closing_box.text_frame.paragraphs[0]
p.text = "From fragmented identities to explainable intelligence."
p.font.size = Pt(16)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.italic = True
p.alignment = PP_ALIGN.CENTER
try:
    closing_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass

# Safety / accuracy disclaimer
add_text(slide6, Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.3),
         "Defensive research only. Attribution reflects analytical correlation & confidence — not definitive identification of real individuals. All data is synthetic.",
         font_size=7.5, color=WHITE_MUTE, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape(slide6, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=CYAN)


# ═══════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════

output_path = "/home/darkbytehunter/Desktop/DarkTrace-Nexus/DarkTraceNexus_SIH_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(output_path) / 1024:.1f} KB")
