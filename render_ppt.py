#!/usr/bin/env python3
"""Render PPTX slides to PNG using PIL for visual inspection."""
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

PX_PER_INCH = 110
SCALE = PX_PER_INCH / 914400.0  # EMU -> px

def emu_px(v):
    return int(v * SCALE)

def get_font(size_pt, bold=False):
    px = max(8, int(size_pt * PX_PER_INCH / 72))
    # Try common fonts
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    for c in candidates:
        try:
            return ImageFont.truetype(c, px)
        except:
            pass
    return ImageFont.load_default()

def rgb_of(color):
    if color is None:
        return None
    try:
        return (color[0], color[1], color[2])
    except:
        return None

def draw_text(draw, shape, im):
    tf = shape.text_frame
    # bounding box
    x0 = emu_px(shape.left)
    y0 = emu_px(shape.top)
    x1 = emu_px(shape.left + shape.width)
    y1 = emu_px(shape.top + shape.height)
    # vertical anchor
    anchor = 'top'
    try:
        anchor = str(tf.vertical_anchor)
    except:
        pass
    # collect lines
    lines = []
    for p in tf.paragraphs:
        txt = p.text
        size = 12
        color = (230,230,230)
        bold = False
        align = 'left'
        try:
            if p.font.size: size = p.font.size.pt
            if p.font.color and p.font.color.type is not None:
                c = rgb_of(p.font.color.rgb)
                if c: color = c
            bold = bool(p.font.bold)
            align = {PP_ALIGN.CENTER:'center', PP_ALIGN.RIGHT:'right'}.get(p.alignment, 'left')
        except:
            pass
        lines.append((txt, size, color, bold, align))
    # compute total height
    line_h = []
    for (txt, size, color, bold, align) in lines:
        f = get_font(size, bold)
        lh = f.getbbox("Ag")[3] + 4
        line_h.append(lh)
    total = sum(line_h)
    pad = 4
    if 'MIDDLE' in anchor:
        cy = (y0 + y1) / 2 - total/2
    elif 'BOTTOM' in anchor:
        cy = y1 - total - pad
    else:
        cy = y0 + pad
    cx = x0 + pad
    for (txt, size, color, bold, align), lh in zip(lines, line_h):
        f = get_font(size, bold)
        bbox = draw.textbbox((0,0), txt, font=f)
        tw = bbox[2]-bbox[0]
        if align == 'center':
            tx = (x0 + x1)/2 - tw/2
        elif align == 'right':
            tx = x1 - tw - pad
        else:
            tx = cx
        draw.text((tx, cy), txt, fill=color, font=f)
        cy += lh

def draw_shape(draw, shape):
    st = None
    try:
        st = shape.auto_shape_type
    except:
        st = None
    stype = str(st) if st else ""
    x0 = emu_px(shape.left)
    y0 = emu_px(shape.top)
    x1 = emu_px(shape.left + shape.width)
    y1 = emu_px(shape.top + shape.height)
    # fill
    fill = None
    try:
        if shape.fill.type is not None and shape.fill.fore_color and shape.fill.fore_color.type is not None:
            fill = rgb_of(shape.fill.fore_color.rgb)
    except:
        fill = None
    # line
    line = None
    try:
        if shape.line.color and shape.line.color.type is not None:
            line = rgb_of(shape.line.color.rgb)
    except:
        line = None
    line_w = 1
    try:
        if shape.line.width: line_w = max(1, int(shape.line.width/914400*PX_PER_INCH))
    except:
        pass

    if 'ARROW' in stype.upper():
        # draw arrow as filled triangle pointing direction
        color = fill or (0,212,255)
        cx, cy = (x0+x1)/2, (y0+y1)/2
        w, h = x1-x0, y1-y0
        if 'RIGHT' in stype.upper():
            pts = [(x0, y0), (x1, cy), (x0, y1)]
        elif 'LEFT' in stype.upper():
            pts = [(x1, y0), (x0, cy), (x1, y1)]
        elif 'DOWN' in stype.upper():
            pts = [(x0, y0), (cx, y1), (x1, y0)]
        elif 'UP' in stype.upper():
            pts = [(x0, y1), (cx, y0), (x1, y1)]
        else:
            pts = [(x0,y0),(x1,y1)]
        draw.polygon(pts, fill=color)
        return
    if 'OVAL' in stype.upper() or stype == 'OVAL':
        draw.ellipse([x0,y0,x1,y1], fill=fill, outline=line, width=line_w)
        return
    # default rounded rect / rect
    draw.rounded_rectangle([x0,y0,x1,y1], radius=6, fill=fill, outline=line, width=line_w)

def render(prs, idx, out):
    slide = prs.slides[idx]
    W = emu_px(prs.slide_width)
    H = emu_px(prs.slide_height)
    im = Image.new("RGB", (W, H), (10,14,26))
    d = ImageDraw.Draw(im)
    # background
    try:
        bg = slide.background.fill
        if bg.fore_color and bg.fore_color.type is not None:
            c = rgb_of(bg.fore_color.rgb)
            if c: im = Image.new("RGB", (W,H), c); d = ImageDraw.Draw(im)
    except:
        pass
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (shape.has_text_frame and shape.text_frame.text.strip()):
                # draw fill if any then text
                if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                    draw_shape(d, shape)
                draw_text(d, shape, im)
            else:
                draw_shape(d, shape)
        except Exception as e:
            print(f"  shape err: {e}")
    im.save(out)
    print(f"  saved {out} ({W}x{H})")

if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv)>1 else "DarkTraceNexus_SIH_Presentation.pptx"
    prs = Presentation(path)
    import os
    os.makedirs("/tmp/pptrender", exist_ok=True)
    only = [int(x) for x in sys.argv[2:]] if len(sys.argv)>2 else range(len(prs.slides))
    for i in only:
        render(prs, i, f"/tmp/pptrender/slide_{i+1}.png")
